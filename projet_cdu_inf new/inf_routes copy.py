from flask import Blueprint, send_file
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt, Cm
from pptx.enum.chart import XL_LABEL_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from db import get_connection
from io import BytesIO
import copy
import re

inf_bp = Blueprint('inf', __name__, url_prefix='/inf')

def replace_placeholder(slide, replacements: dict):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for placeholder, replacement in replacements.items():
                    if placeholder in run.text:
                        run.text = run.text.replace(placeholder, replacement)

def clean_observ(text):
    if not text:
        return ""
    return re.sub(r'_x000D_|_x000A_|\r|\n', '\n', text).strip()

def estimate_text_height(text, font_size_pt=14, line_spacing_pt=16, max_height_cm=12):
    lines = text.count('\n') + 1
    total_height_pt = lines * line_spacing_pt
    cm_per_pt = 0.0352778
    total_height_cm = total_height_pt * cm_per_pt
    return total_height_cm <= max_height_cm

def split_by_estimated_height(text, max_height_cm=12):
    paragraphs = text.strip().split('\n\n')
    chunks = []
    current_chunk = ""

    for para in paragraphs:
        temp_chunk = current_chunk + para + '\n\n'
        if estimate_text_height(temp_chunk, max_height_cm=max_height_cm):
            current_chunk = temp_chunk
        else:
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
            current_chunk = para + '\n\n'

    if current_chunk.strip():
        chunks.append(current_chunk.strip())
    return chunks

def insert_histogram(slide, categories, values, title_text, x_cm, y_cm, width_cm, height_cm, all_structures):
    full_data = {struct: 0 for struct in all_structures}
    for struct, val in zip(categories, values):
        full_data[struct] = val
    ordered_categories = sorted(full_data.keys())
    ordered_values = [full_data[struct] for struct in ordered_categories]

    chart_data = CategoryChartData()
    chart_data.categories = ordered_categories
    chart_data.add_series('Valeurs', ordered_values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Cm(x_cm), Cm(y_cm), Cm(width_cm), Cm(height_cm), chart_data
    ).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = title_text

    for series in chart.series:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = RGBColor(255, 165, 0)
        series.has_data_labels = True
        labels = series.data_labels
        labels.show_value = True
        labels.position = XL_LABEL_POSITION.OUTSIDE_END
        labels.font.size = Pt(10)
        labels.font.bold = True

    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.category_axis.tick_labels.rotation = -45
    chart.value_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.bold = True

def add_total_text(slide, total_value, x_cm, y_cm, width_cm, height_cm):
    textbox = slide.shapes.add_textbox(Cm(x_cm), Cm(y_cm), Cm(width_cm), Cm(height_cm))
    p = textbox.text_frame.paragraphs[0]
    p.text = f"Total : {total_value}"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.CENTER

@inf_bp.route('/generate', methods=['GET'])
def generate_pptx():
    try:
        conn = get_connection()
        cur = conn.cursor()

        def safe_fetch(query):
            cur.execute(query)
            return cur.fetchall() or []

        ged_mois = safe_fetch("SELECT STRUCT, SUM(NBR_INSERTION) FROM CDUGL1Z.GED_MOIS GROUP BY STRUCT")
        ged_total = safe_fetch("SELECT STRUCT, SUM(NBR_INSERTION) FROM CDUGL1Z.GED_TOTAL GROUP BY STRUCT")
        inf = safe_fetch("""SELECT STRUCT, SUM(VALEUR) FROM CDUGL1Z.INF 
                             WHERE TYPE_INTERVENTION IN ('SUPPORT','Support','BDD','RESEAUX') 
                             GROUP BY STRUCT""")
        tel = safe_fetch("""SELECT STRUCT, SUM(VALEUR) FROM CDUGL1Z.INF 
                             WHERE TYPE_INTERVENTION = 'TELECOM' 
                             GROUP BY STRUCT""")

        all_structures = sorted(set([s for s, _ in ged_mois + ged_total + inf + tel]))

        dict_ged = dict(ged_mois)
        dict_ged_total = dict(ged_total)
        dict_inf = dict(inf)
        dict_tel = dict(tel)

        valeurs_ged = [dict_ged.get(s, 0) for s in all_structures]
        valeurs_ged_total = [dict_ged_total.get(s, 0) for s in all_structures]
        valeurs_inf = [dict_inf.get(s, 0) for s in all_structures]
        valeurs_tel = [dict_tel.get(s, 0) for s in all_structures]

        total_inf = sum(valeurs_inf)
        total_tel = sum(valeurs_tel)

        observations = safe_fetch("""SELECT TYPE_INTERVENTION, OBSERV FROM CDUGL1Z.INF 
                                      WHERE TYPE_INTERVENTION IN ('TELECOM', 'RESEAUX', 'BDD') 
                                      AND OBSERV IS NOT NULL""")

        cur.execute("SELECT MAX(PERIODE) FROM CDUGL1Z.INF")
        row = cur.fetchone()
        periode = str(row[0]) if row and row[0] else "Période inconnue"

        cur.close()
        conn.close()

        prs = Presentation("PPTX/SIG MAI template.pptx")

        insert_histogram(prs.slides[2], all_structures, valeurs_ged, f"GED Effort {periode}", 6.56, 0.03, 24.35, 4.76, all_structures)
        insert_histogram(prs.slides[2], all_structures, valeurs_ged_total, "GED TOTAL", 6.56, 5.33, 24.35, 4.76, all_structures)
        insert_histogram(prs.slides[3], all_structures, valeurs_inf, "Interventions INF", 1.86, 3.46, 28.3, 11.37, all_structures)
        insert_histogram(prs.slides[4], all_structures, valeurs_tel, "Interventions TELECOM", 1.86, 3.46, 28.3, 11.37, all_structures)

        add_total_text(prs.slides[3], total_inf, 10, 15.2, 10, 1)
        add_total_text(prs.slides[4], total_tel, 10, 15.2, 10, 1)

        for slide in prs.slides:
            replace_placeholder(slide, {"{{PERIODE}}": periode})

        model_slide = prs.slides[5]

        for ttype, observ in observations:
            clean_text = clean_observ(observ)
            if not clean_text:
                continue

            chunks = split_by_estimated_height(clean_text, max_height_cm=11.5)

            for i, chunk in enumerate(chunks):
                new_slide = prs.slides.add_slide(model_slide.slide_layout)
                for shape in model_slide.shapes:
                    el = copy.deepcopy(shape.element)
                    new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')

                suffix = f" (suite {i})" if i > 0 else ""
                replace_placeholder(new_slide, {
                    "{{PERIODE}}": periode,
                    "{{TYPE_INTERVENTION}}": f"{ttype}{suffix}",
                    "{{OBSERV}}": chunk
                })

        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        if len(slides) > 5:
            xml_slides.remove(slides[5])

        pptx_io = BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)

        return send_file(
            pptx_io,
            as_attachment=True,
            download_name="rapport_SIG_final.pptx",
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    except Exception as e:
        return f"Erreur lors de la génération : {e}", 500
